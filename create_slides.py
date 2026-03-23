#!/usr/bin/env python3
"""Generate KubeCon Amsterdam presentation slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Color palette
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)       # Dark navy
ACCENT_BLUE = RGBColor(0x00, 0x9A, 0xF0)    # Bright blue
ACCENT_GREEN = RGBColor(0x00, 0xC8, 0x53)   # Green
ACCENT_RED = RGBColor(0xFF, 0x45, 0x45)      # Red
ACCENT_YELLOW = RGBColor(0xFF, 0xC1, 0x07)   # Yellow
ACCENT_PURPLE = RGBColor(0xBB, 0x86, 0xFC)   # Purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)
MID_GRAY = RGBColor(0x80, 0x80, 0x90)
DARK_CARD = RGBColor(0x25, 0x25, 0x40)

TOTAL_SLIDES = 11


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=LIGHT_GRAY, bullet_color=ACCENT_BLUE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing

        # Bullet character
        run_bullet = p.add_run()
        run_bullet.text = "\u25B8 "  # triangle bullet
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = "Calibri"

        # Item text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = "Calibri"

    return txBox


def add_section_header(slide, text):
    """Add a colored accent bar + section title at the top."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.8),
                 text, font_size=32, color=WHITE, bold=True)


def add_card(slide, left, top, width, height, title, items, title_color=ACCENT_BLUE):
    """Add a card with a title and bullet items."""
    add_shape(slide, left, top, width, height, DARK_CARD,
              border_color=RGBColor(0x35, 0x35, 0x55))

    add_text_box(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.5),
                 title, font_size=20, color=title_color, bold=True)

    add_bullet_list(slide, left + Inches(0.2), top + Inches(0.6), width - Inches(0.4),
                    height - Inches(0.75), items, font_size=18, color=LIGHT_GRAY,
                    bullet_color=title_color, spacing=Pt(6))


def add_speaker_notes(slide, text):
    """Add speaker notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_slide_number(slide, num):
    """Add slide number at bottom-right."""
    add_text_box(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num} / {TOTAL_SLIDES}", font_size=11, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_arrow(slide, left, top, width, height, color, label=""):
    """Add an arrow shape with optional label."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    if label:
        add_text_box(slide, left, top - Inches(0.25), width, Inches(0.25),
                     label, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    return arrow


# ============================================================
# Create presentation (16:9 widescreen)
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BG)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
             "Signed, Sealed, Delivered",
             font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
             "Identity-Based Access Control with Reverse Proxies",
             font_size=28, color=ACCENT_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(5), Inches(3.8), Inches(3.33), Pt(2))
divider.fill.solid()
divider.fill.fore_color.rgb = ACCENT_BLUE
divider.line.fill.background()

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
             "KubeCon + CloudNativeCon Europe 2026  |  Amsterdam",
             font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Speaker names
add_text_box(slide, Inches(1), Inches(4.9), Inches(11), Inches(0.5),
             "Peter ONeill & Boris Kurktchiev",
             font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Keycloak  \u00B7  Envoy Proxy  \u00B7  OAuth2/OIDC  \u00B7  JWT  \u00B7  Docker",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.4), Inches(13.33), Pt(6))
bar2.fill.solid()
bar2.fill.fore_color.rgb = ACCENT_BLUE
bar2.line.fill.background()

add_slide_number(slide, 1)
add_speaker_notes(slide, """Welcome to "Signed, Sealed, Delivered: Identity-Based Access Control with Reverse Proxies"

Key setup:
- This talk demonstrates how to move beyond VPN-based access control
- We'll show a live demo using Keycloak, Envoy, and JWT tokens
- The core insight: authentication and authorization are different things
- By the end, you'll see how a reverse proxy can enforce per-user, per-resource access control with full audit trails

Introduce yourselves briefly, then move to speaker intro slide.""")


# ============================================================
# SLIDE 2: Speaker Intro
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "About Us")

card_w = Inches(5.5)
card_h = Inches(3.5)
peter_card = add_shape(slide, Inches(0.8), Inches(1.8), card_w, card_h, DARK_CARD,
                       border_color=ACCENT_BLUE)

add_text_box(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.6),
             "Peter ONeill", font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(0.4),
             "Teleport", font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(3.2), Inches(4.8), Inches(1.8),
                ["Solutions Engineering",
                 "Infrastructure & Zero Trust Security",
                 "github.com/peteroneilljr"],
                font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT_BLUE)

boris_card = add_shape(slide, Inches(7.0), Inches(1.8), card_w, card_h, DARK_CARD,
                       border_color=ACCENT_PURPLE)

add_text_box(slide, Inches(7.4), Inches(2.0), Inches(4.8), Inches(0.6),
             "Boris Kurktchiev", font_size=28, color=WHITE, bold=True)
add_text_box(slide, Inches(7.4), Inches(2.6), Inches(4.8), Inches(0.4),
             "Teleport", font_size=18, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(3.2), Inches(4.8), Inches(1.8),
                ["Field CTO",
                 "Cloud-Native Plumber",
                 "Security & Identity"],
                font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT_PURPLE)

add_slide_number(slide, 2)
add_speaker_notes(slide, """Speaker introductions:

Peter ONeill:
- Solutions Engineering at Teleport
- Focus on infrastructure and zero trust security
- Built this demo to show practical identity-based access control
- GitHub: github.com/peteroneilljr

Boris Kurktchiev:
- Field CTO at Teleport
- Self-described "Cloud-Native Plumber" - connecting all the pipes
- Deep expertise in security and identity systems
- Brings the architectural perspective on why this matters at scale""")


# ============================================================
# SLIDE 3: Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "What We're Talking About")

add_card(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.5),
         "PROBLEM: The VPN Approach", [
             "Once inside, reach everything",
             "No per-request identity checks",
             "Audit trail limited to IP addresses",
             "Lateral movement is trivial",
         ], title_color=ACCENT_RED)

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.5),
         "SOLUTION: Reverse Proxy + JWT", [
             "Every request carries signed identity",
             "Per-user, per-resource authorization",
             "Backends never handle auth code",
             "Full identity-aware audit trail",
         ], title_color=ACCENT_GREEN)

add_card(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
         "DEMO: What We'll Show", [
             "Keycloak issuing JWTs with identity claims",
             "Envoy validating tokens and enforcing RBAC",
             "Bob (admin) gets 403 on Alice's app",
             "Identity-aware access logs: who, what, when, result",
         ], title_color=ACCENT_BLUE)

add_slide_number(slide, 3)
add_speaker_notes(slide, """Overview of the problem and our solution:

THE VPN PROBLEM:
- Traditional VPNs give network-level access: once you're in, you can reach everything
- No per-request identity verification - just IP-based trust
- Audit trail is limited to IP addresses - "who" is just a network address
- Lateral movement is trivial once inside the network perimeter
- It's all-or-nothing: coarse-grained access control

THE REVERSE PROXY SOLUTION:
- Every HTTP request carries a signed JWT with the user's identity
- The proxy validates identity cryptographically on every request
- Authorization is per-user and per-resource (fine-grained)
- Backend services don't need any authentication code - the proxy handles it
- Complete identity-aware audit trail: who did what, when, with what result

WHAT WE'LL DEMONSTRATE:
- Keycloak as our OAuth2/OIDC identity provider, issuing RS256-signed JWTs
- Envoy as the reverse proxy, validating tokens and enforcing RBAC policies
- The key moment: Bob has an admin role but gets 403 Forbidden on Alice's app
  This proves authentication (who you are) is different from authorization (what you can do)
- Identity-aware access logs that show the full picture for compliance and forensics""")


# ============================================================
# SLIDE 4: OAuth2 / OIDC (simplified)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "OAuth2 & OpenID Connect (OIDC)")

add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.5),
         "The Relationship", [
             "OIDC is built on top of OAuth2",
             "OAuth2 = Authorization (what can you access?)",
             "OIDC = Authentication (who are you?)",
             "Together: identity + permissions in one token",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.5),
         "Components in Our Demo", [
             "Keycloak: Authorization Server + IdP",
             "Envoy: Resource Server + Token Validator",
             "JWT: carries identity + access claims",
             "JWKS Endpoint: public keys for verification",
         ], title_color=ACCENT_PURPLE)

add_card(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(2.5),
         "Demo Flow (Password Grant)", [
             "User sends credentials to Keycloak, gets RS256-signed JWT",
             "JWT contains: username, email, roles, expiration",
             "Client sends JWT as Bearer token to Envoy",
             "Envoy fetches JWKS, validates signature + claims",
         ], title_color=ACCENT_YELLOW)

add_slide_number(slide, 4)
add_speaker_notes(slide, """OAuth2 & OIDC deep dive:

THE RELATIONSHIP:
- OAuth2 is an authorization framework - it answers "what can this client access?"
- OIDC (OpenID Connect) is an authentication layer built ON TOP of OAuth2
- OIDC adds: ID tokens, standard identity claims, discovery endpoints
- Together they give us both identity (who) and permissions (what) in a single token flow

COMPONENTS IN OUR DEMO:
- Keycloak serves as both the Authorization Server (OAuth2) and Identity Provider (OIDC)
- Envoy acts as the Resource Server - it validates tokens and enforces access policies
- The JWT token carries both identity claims (preferred_username, email) and access claims (realm_access.roles)
- The JWKS (JSON Web Key Set) endpoint exposes Keycloak's public keys so Envoy can verify signatures without sharing secrets

DEMO FLOW (6 steps in detail):
1. User POSTs username/password to Keycloak's token endpoint (password grant)
2. Keycloak validates credentials against its user database
3. Keycloak generates an RS256-signed JWT with OIDC claims (iss, sub, aud, exp, preferred_username, email, realm_access.roles)
4. Client receives the JWT and sends it as a Bearer token in the Authorization header to Envoy
5. Envoy fetches Keycloak's public key from the JWKS endpoint, validates the signature and claims (issuer, expiration, audience)
6. If valid, Envoy extracts claims into metadata and forwards the request to the backend

PRODUCTION NOTE: We use Password Grant for demo simplicity. In production, always use Authorization Code Flow + PKCE - never expose user passwords to the client application.""")


# ============================================================
# SLIDE 5: JWT (simplified, split cards)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "JSON Web Tokens (JWT)")

# Header card
add_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.2),
         "Header", [
             "Algorithm: RS256",
             "Type: JWT",
             "Key ID: identifies signing key",
         ], title_color=ACCENT_BLUE)

# Payload card
add_card(slide, Inches(4.7), Inches(1.5), Inches(4.2), Inches(2.2),
         "Payload (Claims)", [
             "iss, sub, aud, exp (standard)",
             "preferred_username: alice",
             "realm_access.roles: [user]",
         ], title_color=ACCENT_GREEN)

# Signature card
add_card(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(2.2),
         "Signature", [
             "RSA-SHA256 signed",
             "Only Keycloak can create",
             "Anyone can verify (JWKS)",
         ], title_color=ACCENT_PURPLE)

# Security properties
add_card(slide, Inches(0.5), Inches(4.2), Inches(6), Inches(2.8),
         "Security Properties", [
             "Integrity: any change breaks signature",
             "Stateless: no DB lookup per request",
             "Time-limited: auto-expires (5 min TTL)",
             "Auditable: identity in every token",
         ], title_color=ACCENT_YELLOW)

# Trade-offs
add_card(slide, Inches(7.0), Inches(4.2), Inches(5.8), Inches(2.8),
         "Trade-offs", [
             "Cannot revoke before expiry (use short TTL)",
             "Payload is encoded, NOT encrypted",
             "Stolen token valid until expiry (use HTTPS)",
         ], title_color=ACCENT_RED)

add_slide_number(slide, 5)
add_speaker_notes(slide, """JWT deep dive:

STRUCTURE - A JWT has three Base64URL-encoded parts separated by dots:

HEADER:
- alg: RS256 (RSA with SHA-256) - asymmetric algorithm, so Envoy only needs the public key
- typ: JWT - identifies this as a JSON Web Token
- kid: Key ID - tells the validator which key from the JWKS to use for verification

PAYLOAD (CLAIMS):
- Standard claims: iss (issuer - Keycloak URL), sub (subject - user UUID), aud (audience), exp (expiration timestamp)
- OIDC claims: preferred_username ("alice"), email, email_verified
- Custom claims: realm_access.roles (["user"] or ["admin"]) - these drive our RBAC policies
- iat (issued at), nbf (not before) - time-based validity bounds

SIGNATURE:
- Created by: Base64URL(header) + "." + Base64URL(payload), signed with Keycloak's RSA private key
- Verification: Envoy fetches the public key from JWKS endpoint and verifies the signature
- This proves the token was issued by Keycloak and hasn't been tampered with

ENVOY VALIDATION (6 steps):
1. Check token format (must have exactly 3 dot-separated parts)
2. Decode header to get algorithm and key ID
3. Fetch matching public key from Keycloak's JWKS endpoint
4. Verify RSA-SHA256 signature using the public key
5. Validate claims: check issuer matches expected, token not expired, audience correct
6. Store validated claims in Envoy's Dynamic Metadata for use by RBAC filter and access logger

TRADE-OFFS:
- Can't revoke a JWT before expiration - mitigate with short TTL (5 min) and refresh tokens
- Payload is Base64URL encoded, NOT encrypted - never put secrets in a JWT
- A stolen token is valid until it expires - always use HTTPS in production""")


# ============================================================
# SLIDE 6: Reverse Proxy (simplified)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "Reverse Proxy Architecture")

add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.2),
         "What is a Reverse Proxy?", [
             "Sits in front of backends, not clients",
             "Single entry point for all traffic",
             "Centralizes auth, authz, logging, TLS",
             "Like a security guard at the entrance",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.2),
         "VPN vs Reverse Proxy", [
             "VPN: network-level, all-or-nothing",
             "Proxy: app-level, per-user, per-resource",
             "VPN: no per-request validation",
             "Proxy: every request verified (zero trust)",
         ], title_color=ACCENT_PURPLE)

add_card(slide, Inches(0.5), Inches(4.2), Inches(7.5), Inches(2.8),
         "Envoy's Filter Chain", [
             "JWT Auth Filter: validates signature, extracts claims",
             "RBAC Filter: checks username vs route policy",
             "Access Logger: logs user + path + decision",
             "Router: forwards to correct backend",
         ], title_color=ACCENT_YELLOW)

add_card(slide, Inches(8.5), Inches(4.2), Inches(4.3), Inches(2.8),
         "Defense in Depth", [
             "Network Isolation (Docker)",
             "JWT Authentication",
             "RBAC Authorization",
             "Identity-Aware Logging",
         ], title_color=ACCENT_GREEN)

add_slide_number(slide, 6)
add_speaker_notes(slide, """Reverse Proxy architecture:

WHAT IS A REVERSE PROXY?
- A reverse proxy sits in front of your backend services - it's the single entry point
- All client traffic flows through it - clients never talk directly to backends
- It centralizes cross-cutting concerns: authentication, authorization, logging, TLS termination
- Think of it like a security guard at a building entrance - everyone must pass through and show ID
- Key difference from a forward proxy: a forward proxy sits in front of clients (like a corporate proxy), a reverse proxy sits in front of servers

VPN vs REVERSE PROXY:
- VPN gives network-level access: once you're connected, you can reach everything on that network
- Reverse proxy gives app-level access: each request is individually authorized based on user identity
- VPN has no per-request validation - once the tunnel is up, all traffic flows through
- Reverse proxy validates every single request with cryptographic identity (zero trust model)
- VPN logs show IP addresses; reverse proxy logs show actual user identities

ENVOY'S FILTER CHAIN (request pipeline):
1. Listener (0.0.0.0:8080) accepts incoming connections
2. JWT Authentication Filter: validates the Bearer token's RS256 signature using JWKS, extracts claims into Dynamic Metadata
3. RBAC Authorization Filter: reads preferred_username from metadata, checks it against per-route policies (e.g., only "alice" can access /alice)
4. Access Logger: writes identity-aware log entries using Dynamic Metadata format strings
5. Router: forwards the validated, authorized request to the correct backend cluster

Invalid token? Rejected in ~5ms. The backend is never even contacted.

DEFENSE IN DEPTH (5 security layers):
1. Network Isolation (Docker network) - backends not directly reachable
2. TLS Termination (in production) - encrypted in transit
3. JWT Authentication - cryptographic identity verification
4. RBAC Authorization - per-user, per-resource access control
5. Rate Limiting (in production) - prevents abuse""")


# ============================================================
# SLIDE 7: Access Logging (simplified)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "Identity-Aware Access Logging")

add_card(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(2.2),
         "How Identity Gets Into Logs", [
             "JWT filter extracts claims into metadata",
             "Access logger reads metadata via format strings",
             "Every log entry includes user identity",
             "One config line makes it work",
         ], title_color=ACCENT_BLUE)

add_card(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.2),
         "The Key Config Line", [
             'payload_in_metadata: "jwt_payload"',
             "Without it: logs show IP only (like VPN)",
             "With it: full identity-aware audit trail",
         ], title_color=ACCENT_YELLOW)

# Three log scenarios
add_card(slide, Inches(0.5), Inches(4.2), Inches(3.8), Inches(2.8),
         "401: No / Invalid Token", [
             'user: "-" (no identity)',
             "JWT rejected before reaching backend",
             "response_flags: UAEX",
         ], title_color=ACCENT_RED)

add_card(slide, Inches(4.7), Inches(4.2), Inches(3.8), Inches(2.8),
         "403: Valid Token, Wrong User", [
             'user: "bob" (identity IS logged)',
             "JWT valid, but RBAC denied access",
             "We know WHO tried to access WHAT",
         ], title_color=ACCENT_YELLOW)

add_card(slide, Inches(8.9), Inches(4.2), Inches(3.9), Inches(2.8),
         "200: Authorized Access", [
             'user: "alice", path: "/alice"',
             "Full audit: who, what, when, result",
             "Compliance: GDPR, SOC2, HIPAA ready",
         ], title_color=ACCENT_GREEN)

add_slide_number(slide, 7)
add_speaker_notes(slide, """Identity-Aware Access Logging:

HOW IDENTITY GETS INTO LOGS:
1. The JWT Authentication filter validates the token and extracts all claims
2. The critical config line: payload_in_metadata: "jwt_payload" - this stores the JWT claims in Envoy's Dynamic Metadata
3. Dynamic Metadata is a request-scoped key-value store that flows through the filter chain
4. The Access Logger reads metadata using format strings like: %DYNAMIC_METADATA(envoy.filters.http.jwt_authn:jwt_payload:preferred_username)%
5. Result: every log entry contains the actual user identity, not just an IP address

THE KEY CONFIG LINE:
- Without payload_in_metadata: logs show only IP addresses, just like a VPN
- With payload_in_metadata: you get full identity-aware audit trails
- This single line is what transforms your logging from network-level to identity-level

THREE SCENARIOS:

401 - No Token or Invalid Token:
- User field shows "-" because the JWT was rejected before metadata could be populated
- The backend was never contacted - fail-fast security
- Response flags show UAEX (Unauthorized External)
- Example: someone tries to access without a token, or with an expired/forged token

403 - Valid Token, Wrong User:
- This is the most interesting case for security
- The user's identity IS logged even though access was denied
- Example: Bob (with admin role) tries to access /alice and gets 403
- Response flags show RBAC_ACCESS_DENIED
- We know exactly WHO tried to access WHAT and WHEN - this is crucial for security monitoring

200 - Authorized Access:
- Complete audit trail: user identity, requested path, timestamp, response code
- Non-repudiation: the user cannot deny performing the action
- Compliance-ready for GDPR (data access tracking), HIPAA (PHI access logging), SOC2 (access audit), PCI-DSS

In production: feed these logs to ELK, Splunk, Datadog, or CloudWatch for alerting and analysis.""")


# ============================================================
# SLIDE 8: Demo Architecture (with arrows)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "Demo Architecture")

# Docker network boundary
docker_border = add_shape(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
                          RGBColor(0x20, 0x20, 0x38), border_color=MID_GRAY)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(4), Inches(0.4),
             "Docker Network (demo-network)", font_size=14, color=MID_GRAY)

# Client box
client = add_shape(slide, Inches(0.8), Inches(2.5), Inches(2), Inches(1), DARK_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(0.9), Inches(2.6), Inches(1.8), Inches(0.4),
             "Client (curl)", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.9), Inches(3.0), Inches(1.8), Inches(0.3),
             "Sends Bearer token", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Keycloak box
kc = add_shape(slide, Inches(0.8), Inches(4.8), Inches(2), Inches(1.5), DARK_CARD, ACCENT_PURPLE)
add_text_box(slide, Inches(0.9), Inches(4.9), Inches(1.8), Inches(0.4),
             "Keycloak", font_size=16, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.9), Inches(5.3), Inches(1.8), Inches(0.8),
             "Port: 8180\nOIDC Provider\nJWKS Endpoint",
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Envoy box (large, center)
envoy = add_shape(slide, Inches(3.8), Inches(2.2), Inches(5), Inches(4.2), DARK_CARD, ACCENT_BLUE)
add_text_box(slide, Inches(4.0), Inches(2.3), Inches(4.6), Inches(0.5),
             "Envoy Proxy  (Port: 8080)", font_size=18, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)

# Filter boxes inside Envoy
f1 = add_shape(slide, Inches(4.1), Inches(2.9), Inches(4.5), Inches(0.7),
               RGBColor(0x30, 0x30, 0x50), ACCENT_YELLOW)
add_text_box(slide, Inches(4.3), Inches(3.0), Inches(4.1), Inches(0.5),
             "1. JWT Authentication Filter", font_size=14, color=ACCENT_YELLOW)

f2 = add_shape(slide, Inches(4.1), Inches(3.7), Inches(4.5), Inches(0.7),
               RGBColor(0x30, 0x30, 0x50), ACCENT_GREEN)
add_text_box(slide, Inches(4.3), Inches(3.8), Inches(4.1), Inches(0.5),
             "2. RBAC Authorization Filter", font_size=14, color=ACCENT_GREEN)

f3 = add_shape(slide, Inches(4.1), Inches(4.5), Inches(4.5), Inches(0.7),
               RGBColor(0x30, 0x30, 0x50), ACCENT_PURPLE)
add_text_box(slide, Inches(4.3), Inches(4.6), Inches(4.1), Inches(0.5),
             "3. Access Logger", font_size=14, color=ACCENT_PURPLE)

f4 = add_shape(slide, Inches(4.1), Inches(5.3), Inches(4.5), Inches(0.7),
               RGBColor(0x30, 0x30, 0x50), LIGHT_GRAY)
add_text_box(slide, Inches(4.3), Inches(5.4), Inches(4.1), Inches(0.5),
             "4. Router", font_size=14, color=LIGHT_GRAY)

# Backend apps
for i, (name, port, user, color) in enumerate([
    ("public-app", "3000", "Any authed user", ACCENT_BLUE),
    ("alice-app", "3002", "Alice ONLY", ACCENT_GREEN),
    ("bob-app", "3001", "Bob ONLY", ACCENT_YELLOW),
]):
    left = Inches(9.8)
    top = Inches(2.2 + i * 1.6)
    app = add_shape(slide, left, top, Inches(2.7), Inches(1.3), DARK_CARD, color)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.1), Inches(2.4), Inches(0.35),
                 name, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.45), Inches(2.4), Inches(0.25),
                 f"Port: {port}", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.75), Inches(2.4), Inches(0.25),
                 user, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# --- ARROWS ---

# Client -> Envoy (horizontal)
add_arrow(slide, Inches(2.9), Inches(2.85), Inches(0.8), Inches(0.3),
          ACCENT_BLUE, "Bearer token")

# Envoy -> backend apps (horizontal, from router area)
for i in range(3):
    top = Inches(2.65 + i * 1.6)
    add_arrow(slide, Inches(8.7), top, Inches(1.0), Inches(0.25), LIGHT_GRAY)

# Client -> Keycloak (vertical arrow - "get token")
# Use a downward arrow shape between client and keycloak
down_arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(1.5), Inches(3.6), Inches(0.3), Inches(1.1))
down_arrow.fill.solid()
down_arrow.fill.fore_color.rgb = ACCENT_PURPLE
down_arrow.line.fill.background()
add_text_box(slide, Inches(0.1), Inches(3.9), Inches(1.3), Inches(0.3),
             "get token", font_size=14, color=ACCENT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Envoy -> Keycloak (horizontal arrow - "fetch JWKS")
add_arrow(slide, Inches(3.0), Inches(5.35), Inches(0.7), Inches(0.25),
          ACCENT_YELLOW, "fetch JWKS")

add_slide_number(slide, 8)
add_speaker_notes(slide, """Demo Architecture walkthrough:

THE SETUP:
All components run in a single Docker network (demo-network) for isolation.

COMPONENTS:
1. Client (curl) - represents any HTTP client sending Bearer tokens
2. Keycloak (port 8180) - our OIDC provider, issues RS256-signed JWTs, exposes JWKS endpoint
3. Envoy Proxy (port 8080) - the reverse proxy, single entry point
4. Three backend apps: public-app (3000), bob-app (3001), alice-app (3002)

REQUEST FLOW:
1. Client first gets a token from Keycloak (password grant: POST to /token endpoint)
2. Client sends request to Envoy with "Authorization: Bearer <token>" header
3. Envoy's JWT Authentication Filter validates the signature by fetching JWKS from Keycloak
4. Envoy's RBAC Filter checks the preferred_username claim against route-specific policies
5. Access Logger records the user identity + path + decision
6. Router forwards to the correct backend cluster

ENVOY FILTER CHAIN:
- JWT Auth Filter: validates RS256 signature, checks iss/exp/aud, extracts claims into Dynamic Metadata
- RBAC Filter: reads preferred_username from metadata, checks per-route policies
  - /public: any authenticated user
  - /alice: only preferred_username == "alice"
  - /bob: only preferred_username == "bob"
- Access Logger: uses Dynamic Metadata format strings to include identity in every log line
- Router: forwards validated requests to backend clusters with path rewriting

ACCESS CONTROL MATRIX:
- No token → 401 Unauthorized (for all routes)
- Alice → /public: 200, /alice: 200, /bob: 403
- Bob → /public: 200, /alice: 403, /bob: 200
- Key insight: Bob's admin role doesn't help - RBAC checks username, not role""")


# ============================================================
# SLIDE 9: The Aha Moment (NEW)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "The Aha Moment")

# Big callout text
add_text_box(slide, Inches(1), Inches(1.3), Inches(11.33), Inches(0.6),
             "Authentication \u2260 Authorization",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Three large colored boxes showing access
box_w = Inches(3.8)
box_h = Inches(3.5)

# Alice's access
alice_box = add_shape(slide, Inches(0.5), Inches(2.5), box_w, box_h, DARK_CARD,
                      border_color=ACCENT_GREEN)
add_text_box(slide, Inches(0.7), Inches(2.7), Inches(3.4), Inches(0.5),
             "Alice (role: user)", font_size=22, color=ACCENT_GREEN, bold=True,
             alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(0.7), Inches(3.3), Inches(3.4), Inches(2.2),
                ["/public  \u2192  200 OK",
                 "/alice    \u2192  200 OK",
                 "/bob      \u2192  403 Forbidden"],
                font_size=20, color=LIGHT_GRAY, bullet_color=ACCENT_GREEN, spacing=Pt(12))

# Bob's access
bob_box = add_shape(slide, Inches(4.8), Inches(2.5), box_w, box_h, DARK_CARD,
                    border_color=ACCENT_YELLOW)
add_text_box(slide, Inches(5.0), Inches(2.7), Inches(3.4), Inches(0.5),
             "Bob (role: admin)", font_size=22, color=ACCENT_YELLOW, bold=True,
             alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(5.0), Inches(3.3), Inches(3.4), Inches(2.2),
                ["/public  \u2192  200 OK",
                 "/alice    \u2192  403 Forbidden",
                 "/bob      \u2192  200 OK"],
                font_size=20, color=LIGHT_GRAY, bullet_color=ACCENT_YELLOW, spacing=Pt(12))

# Key insight box
insight_box = add_shape(slide, Inches(9.1), Inches(2.5), box_w, box_h, DARK_CARD,
                        border_color=ACCENT_RED)
add_text_box(slide, Inches(9.3), Inches(2.7), Inches(3.4), Inches(0.5),
             "The Key Insight", font_size=22, color=ACCENT_RED, bold=True,
             alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(9.3), Inches(3.3), Inches(3.4), Inches(2.2),
                ["Bob IS authenticated",
                 "Bob has admin role",
                 "Bob still gets 403 on /alice",
                 "RBAC checks identity, not role"],
                font_size=18, color=LIGHT_GRAY, bullet_color=ACCENT_RED, spacing=Pt(10))

# Bottom callout
callout = add_shape(slide, Inches(1.5), Inches(6.3), Inches(10.33), Inches(0.8),
                    ACCENT_BLUE, border_color=ACCENT_BLUE)
add_text_box(slide, Inches(1.7), Inches(6.35), Inches(9.93), Inches(0.7),
             "Valid identity token + elevated role  \u2260  access to everything",
             font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 9)
add_speaker_notes(slide, """THE AHA MOMENT - This is the core takeaway of the entire talk.

SETUP:
- Alice has the role "user" in Keycloak
- Bob has the role "admin" in Keycloak
- Both have valid, RS256-signed JWT tokens from the same identity provider

ALICE'S ACCESS:
- /public → 200 OK (any authenticated user can access)
- /alice → 200 OK (RBAC policy: preferred_username must be "alice")
- /bob → 403 Forbidden (RBAC policy: preferred_username must be "bob")

BOB'S ACCESS:
- /public → 200 OK (any authenticated user can access)
- /alice → 403 FORBIDDEN (RBAC policy: preferred_username must be "alice" - Bob's "admin" role doesn't matter)
- /bob → 200 OK (RBAC policy: preferred_username must be "bob")

THE KEY INSIGHT:
Bob IS authenticated. His JWT is perfectly valid. He even has the "admin" role.
But he still gets 403 Forbidden on /alice's app.

WHY? Because our RBAC policy checks the USERNAME, not the role.
Authentication (proving who you are) is fundamentally different from authorization (what you're allowed to do).

This is zero trust in action:
- Having a valid identity token doesn't grant blanket access
- Even elevated roles don't bypass resource-level authorization
- Every resource has its own access policy based on identity attributes
- The proxy enforces these policies consistently, backends don't need to know

In the VPN world: if Bob is "inside the network", he can reach everything.
In the reverse proxy world: Bob's valid token + admin role still can't access Alice's resources.

This is the "Signed, Sealed, Delivered" moment - the token is signed (valid), sealed (secure), but access is only delivered to the right identity.""")


# ============================================================
# SLIDE 10: Key Takeaways (reduced to 4)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_section_header(slide, "Key Takeaways")

takeaways = [
    ("Zero Trust > VPN",
     "Every request validated with cryptographic identity. Never trust the network.",
     ACCENT_BLUE),
    ("Auth \u2260 Authz",
     "Valid token + admin role still gets 403. Identity and access are separate concerns.",
     ACCENT_GREEN),
    ("Centralized Security",
     "One proxy handles auth for all backends. Implement security once, consistently.",
     ACCENT_PURPLE),
    ("Identity-Aware Logs",
     "Logs show WHO did WHAT, WHEN, with WHAT RESULT. Compliance and forensics built in.",
     ACCENT_YELLOW),
]

for i, (title, desc, color) in enumerate(takeaways):
    col = i % 2
    row = i // 2
    left = Inches(0.5 + col * 6.3)
    top = Inches(1.5 + row * 2.8)
    width = Inches(5.9)
    height = Inches(2.5)

    card = add_shape(slide, left, top, width, height, DARK_CARD, color)

    num_box = add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.5),
                           str(i + 1), font_size=22, color=color, bold=True)

    add_text_box(slide, left + Inches(0.7), top + Inches(0.15), width - Inches(0.9), Inches(0.5),
                 title, font_size=22, color=WHITE, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.8), width - Inches(0.6), height - Inches(1.0),
                 desc, font_size=18, color=LIGHT_GRAY)

add_slide_number(slide, 10)
add_speaker_notes(slide, """Key Takeaways - Four things to remember:

1. ZERO TRUST > VPN:
- VPNs trust the network perimeter. Once inside, everything is reachable.
- Zero trust means never trust, always verify. Every single request is validated with cryptographic identity.
- The reverse proxy is the enforcement point - it validates the JWT signature, checks claims, and enforces RBAC on every request.
- This isn't just theoretical - we demonstrated it live.

2. AUTHENTICATION ≠ AUTHORIZATION:
- This is the biggest misconception in security.
- Authentication proves WHO you are (valid JWT, correct credentials)
- Authorization decides WHAT you can do (RBAC policies, resource-level access)
- Bob proved: valid token + admin role ≠ access to everything
- Design your systems so these are separate concerns, enforced at different layers.

3. CENTRALIZED SECURITY:
- The reverse proxy handles authentication and authorization for ALL backends
- Backends focus on business logic - they never need auth code
- Security is implemented once, in one place, consistently
- Adding a new backend? Just add a route and RBAC policy. No code changes.
- This eliminates the "every team implements auth differently" problem.

4. IDENTITY-AWARE LOGS:
- Traditional logs show IP addresses. Our logs show actual user identities.
- Every log entry: WHO (username), WHAT (path), WHEN (timestamp), RESULT (status code)
- This is essential for: compliance (GDPR, HIPAA, SOC2), incident response, anomaly detection, forensic analysis
- The 403 case is most valuable: we know WHO tried to access WHAT they shouldn't have.
- One config line (payload_in_metadata) transforms your logging from network-level to identity-level.""")


# ============================================================
# SLIDE 11: Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_BLUE
bar.line.fill.background()

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "Questions & Answers",
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(5.5), Inches(3.3), Inches(2.33), Pt(2))
divider.fill.solid()
divider.fill.fore_color.rgb = ACCENT_BLUE
divider.line.fill.background()

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "github.com/peteroneilljr/kubecon-ams-demo-v2",
             font_size=20, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.4), Inches(11), Inches(0.5),
             "goteleport.com",
             font_size=20, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
             "Try it yourself:  docker-compose up -d  &&  ./demo-script.sh",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "Peter ONeill  \u00B7  Boris Kurktchiev  \u00B7  Teleport",
             font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.4), Inches(13.33), Pt(6))
bar2.fill.solid()
bar2.fill.fore_color.rgb = ACCENT_BLUE
bar2.line.fill.background()

add_slide_number(slide, 11)
add_speaker_notes(slide, """Q&A slide.

Key points to revisit if asked:
- The demo repo is fully open source: github.com/peteroneilljr/kubecon-ams-demo-v2
- Everything runs with docker-compose up -d and the demo script
- Teleport (goteleport.com) implements these same zero-trust principles at enterprise scale
- For production: use Authorization Code Flow + PKCE instead of password grant
- For production: add TLS, rate limiting, and connect logs to a SIEM

Common questions to prepare for:
- "How does this scale?" - Envoy is battle-tested at massive scale (Lyft, Google, etc.)
- "What about token revocation?" - Short TTL (5 min) + refresh tokens. For immediate revocation, consider a token blacklist or Envoy's ext_authz filter
- "Can you use this with service mesh?" - Yes, Envoy is the data plane for Istio. Same JWT validation works for service-to-service auth
- "How do you handle multiple identity providers?" - Envoy supports multiple JWT providers in the same filter config
- "What about API keys?" - You can combine JWT auth with API key auth using Envoy's ext_authz filter""")


# ============================================================
# Save
# ============================================================
output_path = "/Users/darkmatter/GitHub/demos/kubecon-ams-demo-v2/kubecon-ams-2026.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
